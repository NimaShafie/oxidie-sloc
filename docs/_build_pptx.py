"""
Generate docs/oxide-sloc-overview.pptx
Run with: python docs/_build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ────────────────────────────────────────────────────────────────────
OX          = RGBColor(0xB8, 0x5C, 0x1A)   # oxide orange
OX_DARK     = RGBColor(0x7A, 0x3A, 0x0D)   # deep brown
CHARCOAL    = RGBColor(0x1E, 0x1A, 0x14)   # near-black
WARM_WHITE  = RGBColor(0xFA, 0xF8, 0xF5)   # off-white
LIGHT_SAND  = RGBColor(0xF0, 0xE8, 0xDC)   # light sand
MUTED       = RGBColor(0x6B, 0x57, 0x44)   # warm grey
GREEN       = RGBColor(0x2A, 0x68, 0x46)   # success green
DARK_BG     = RGBColor(0x18, 0x14, 0x0E)   # dark slide bg

# ── Slide dimensions (widescreen 16:9) ─────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # totally blank layout

# ─────────────────────────────────────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = line_width
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, text, left, top, width, height,
             font_name="Segoe UI", font_size=Pt(18), bold=False, italic=False,
             color=CHARCOAL, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=Pt(16), bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, space_before=Pt(4), indent_level=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.level = indent_level
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name = "Segoe UI"
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def accent_bar(slide, color=OX, height=Inches(0.08)):
    """Full-width accent bar at the top of a slide."""
    add_rect(slide, 0, 0, W, height, fill_rgb=color)


def slide_number(slide, n, total, color=MUTED):
    """Small slide-number in bottom-right corner."""
    add_text(slide, f"{n} / {total}",
             W - Inches(1.2), H - Inches(0.35), Inches(1.1), Inches(0.28),
             font_size=Pt(9), color=color, align=PP_ALIGN.RIGHT)


def footer_line(slide, text="oxide-sloc · Source-Line Metrics for Engineering Teams"):
    add_rect(slide, 0, H - Inches(0.38), W, Inches(0.38), fill_rgb=CHARCOAL)
    add_text(slide, text, Inches(0.3), H - Inches(0.34), W - Inches(1.5), Inches(0.3),
             font_size=Pt(9), color=RGBColor(0xCC, 0xBB, 0xAA), align=PP_ALIGN.LEFT)


def bullet_block(slide, items, left, top, width, height,
                 font_size=Pt(15), color=CHARCOAL, bullet_color=OX):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        bullet_run = p.add_run()
        bullet_run.text = "•  "
        bullet_run.font.name = "Segoe UI"
        bullet_run.font.size = font_size
        bullet_run.font.color.rgb = bullet_color
        text_run = p.add_run()
        text_run.text = item
        text_run.font.name = "Segoe UI"
        text_run.font.size = font_size
        text_run.font.color.rgb = color
    return txBox


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE DATA
# ─────────────────────────────────────────────────────────────────────────────

slides_data = [
    # 0 - Cover
    {"type": "cover"},
    # 1 - What is SLOC?
    {"type": "content",
     "title": "What is SLOC?",
     "subtitle": "Understanding the fundamental metric",
     "bullets": [
         "SLOC = Source Lines of Code",
         "A count of meaningful lines in a software project's source files",
         "Excludes blank lines and comment-only lines (by default)",
         "Widely used in industry for project sizing, estimation, and tracking",
         "Standardised by IEEE 1045-1992 for consistent, reproducible measurement",
     ],
     "note": "SLOC is not a measure of productivity — it is a measure of codebase size and growth."},
    # 2 - Why Measure SLOC?
    {"type": "two_col",
     "title": "Why Measure SLOC?",
     "col1_title": "For Teams & Managers",
     "col1": [
         "Understand project scope at a glance",
         "Track codebase growth sprint-over-sprint",
         "Identify languages and components growing fastest",
         "Support project estimation and resourcing",
         "Provide evidence for technical-debt discussions",
     ],
     "col2_title": "For Engineers",
     "col2": [
         "Monitor the effect of refactors and cleanup",
         "Ensure test code grows alongside production code",
         "Set and enforce code-size budgets in CI pipelines",
         "Detect unexpectedly large files or modules",
         "Generate auditable metrics for compliance reports",
     ]},
    # 3 - Introducing oxide-sloc
    {"type": "highlight",
     "title": "Introducing oxide-sloc",
     "tagline": "A fast, air-gapped-capable SLOC workbench built for enterprise teams",
     "points": [
         ("60 Languages", "Rust, Python, Java, C/C++, TypeScript, SQL, and 54 more"),
         ("Zero Dependencies", "Runs entirely offline — no internet, no pre-installed Rust required"),
         ("Multiple Surfaces", "CLI for scripts · Web UI for ad-hoc use · REST API for automation"),
         ("Rich Reports", "HTML, PDF, CSV, Excel — beautiful charts included"),
         ("CI-Native", "GitHub Actions, Jenkins, GitLab CI — drop-in integration"),
     ]},
    # 4 - Key Features Overview
    {"type": "feature_grid",
     "title": "Feature Highlights"},
    # 5 - The Web UI
    {"type": "content",
     "title": "Localhost Web Interface",
     "subtitle": "No installation required — open your browser and start scanning",
     "bullets": [
         "Launch with a single command: bash scripts/run.sh",
         "Browse and select any directory using a native folder picker",
         "Results appear instantly with interactive charts and a per-file table",
         "Compare two past scans side-by-side to see what changed",
         "Trend charts show SLOC growth over multiple weeks or sprints",
         "Test Metrics page tracks unit test files and test-to-code ratio",
         "Dark mode, 5 colour themes, and light/dark toggle built in",
     ],
     "note": "The web UI binds to 127.0.0.1 by default — only your machine can reach it."},
    # 6 - Who Benefits?
    {"type": "who_benefits"},
    # 7 - How Often Should You Run It?
    {"type": "cadence"},
    # 8 - Integrations
    {"type": "content",
     "title": "Integrations & Delivery",
     "subtitle": "Get SLOC data wherever your team already works",
     "bullets": [
         "GitHub / GitLab:  auto-comment SLOC diffs on every pull request",
         "Jenkins / CI:  fail builds if codebase exceeds size budgets",
         "Microsoft Teams:  Adaptive Card summary posted to any channel",
         "Email (SMTP):  scheduled HTML report delivered to a distribution list",
         "Confluence:  SLOC page created or updated automatically after each scan",
         "REST API:  /api/metrics/latest for dashboards, Grafana, Power BI",
         "Webhook:  POST JSON to any HTTPS endpoint for custom pipelines",
     ]},
    # 9 - Air-Gapped / Secure Environments
    {"type": "content",
     "title": "Built for Secure Environments",
     "subtitle": "Full functionality with no internet access and no pre-installed tools",
     "bullets": [
         "All ~328 Rust dependencies are committed to the repository",
         "Bundled Rust toolchain (Windows, Linux x64, Linux arm64) — no rustup needed",
         "git clone → bash scripts/run.sh is the only step required",
         "Supports RHEL 8/9, Windows 10/11 (Git Bash), Ubuntu, Debian",
         "Optional native TLS and API-key authentication for server deployments",
         "Docker image available for persistent server installations",
         "IP rate limiting and security headers enabled by default",
     ],
     "note": "Designed for air-gapped labs and corporate environments with strict egress controls."},
    # 10 - Getting Started
    {"type": "getting_started"},
    # 11 - Thank you / Q&A
    {"type": "closing"},
]

TOTAL = len(slides_data)

# ─────────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────────

def build_cover(prs):
    slide = prs.slides.add_slide(BLANK)
    # Dark background
    add_rect(slide, 0, 0, W, H, fill_rgb=DARK_BG)
    # Left accent stripe
    add_rect(slide, 0, 0, Inches(0.18), H, fill_rgb=OX)
    # Bottom bar
    add_rect(slide, 0, H - Inches(1.4), W, Inches(1.4), fill_rgb=RGBColor(0x28, 0x22, 0x18))

    # Title
    add_text(slide, "oxide-sloc",
             Inches(0.55), Inches(1.8), Inches(9), Inches(1.5),
             font_size=Pt(72), bold=True, color=OX, align=PP_ALIGN.LEFT)
    add_text(slide, "Source-Line Metrics for Engineering Teams",
             Inches(0.55), Inches(3.2), Inches(10), Inches(0.7),
             font_size=Pt(26), bold=False, color=WARM_WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "Product Overview",
             Inches(0.55), Inches(3.9), Inches(8), Inches(0.5),
             font_size=Pt(18), color=RGBColor(0xAA, 0x88, 0x66), align=PP_ALIGN.LEFT)

    # Bottom meta
    add_text(slide, "What is SLOC?  ·  Who benefits?  ·  How often to run?  ·  Getting started",
             Inches(0.55), H - Inches(1.1), Inches(11), Inches(0.4),
             font_size=Pt(12), color=RGBColor(0x88, 0x77, 0x66), align=PP_ALIGN.LEFT)
    add_text(slide, "github.com/oxide-sloc/oxide-sloc",
             Inches(0.55), H - Inches(0.7), Inches(8), Inches(0.35),
             font_size=Pt(11), color=OX, align=PP_ALIGN.LEFT)


def build_content(prs, n, data):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=WARM_WHITE)
    accent_bar(slide)
    footer_line(slide)
    slide_number(slide, n, TOTAL)

    # Title area
    add_rect(slide, 0, Inches(0.08), W, Inches(1.15), fill_rgb=LIGHT_SAND)
    add_text(slide, data["title"],
             Inches(0.45), Inches(0.14), Inches(11), Inches(0.6),
             font_size=Pt(30), bold=True, color=OX_DARK)
    if "subtitle" in data:
        add_text(slide, data["subtitle"],
                 Inches(0.45), Inches(0.72), Inches(11), Inches(0.42),
                 font_size=Pt(14), color=MUTED)

    # Bullets
    bullet_block(slide, data["bullets"],
                 Inches(0.5), Inches(1.45), Inches(12.1), Inches(4.6),
                 font_size=Pt(16), color=CHARCOAL)

    if "note" in data:
        add_rect(slide, Inches(0.4), H - Inches(1.05), Inches(12.5), Inches(0.6),
                 fill_rgb=RGBColor(0xFF, 0xF4, 0xE8))
        add_text(slide, f"Note: {data['note']}",
                 Inches(0.55), H - Inches(1.02), Inches(12.2), Inches(0.52),
                 font_size=Pt(11), italic=True, color=OX_DARK)


def build_two_col(prs, n, data):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=WARM_WHITE)
    accent_bar(slide)
    footer_line(slide)
    slide_number(slide, n, TOTAL)

    add_rect(slide, 0, Inches(0.08), W, Inches(0.9), fill_rgb=LIGHT_SAND)
    add_text(slide, data["title"],
             Inches(0.45), Inches(0.14), Inches(11), Inches(0.6),
             font_size=Pt(30), bold=True, color=OX_DARK)

    mid = W / 2
    # Divider
    add_rect(slide, mid - Inches(0.015), Inches(1.1), Inches(0.03), H - Inches(1.6),
             fill_rgb=RGBColor(0xD4, 0xC4, 0xB0))

    # Col 1
    add_rect(slide, Inches(0.3), Inches(1.1), mid - Inches(0.65), Inches(0.44),
             fill_rgb=OX)
    add_text(slide, data["col1_title"],
             Inches(0.45), Inches(1.14), mid - Inches(0.8), Inches(0.36),
             font_size=Pt(14), bold=True, color=WARM_WHITE)
    bullet_block(slide, data["col1"],
                 Inches(0.45), Inches(1.65), mid - Inches(0.8), Inches(4.8),
                 font_size=Pt(14))

    # Col 2
    add_rect(slide, mid + Inches(0.35), Inches(1.1), mid - Inches(0.65), Inches(0.44),
             fill_rgb=OX_DARK)
    add_text(slide, data["col2_title"],
             mid + Inches(0.5), Inches(1.14), mid - Inches(0.8), Inches(0.36),
             font_size=Pt(14), bold=True, color=WARM_WHITE)
    bullet_block(slide, data["col2"],
                 mid + Inches(0.5), Inches(1.65), mid - Inches(0.8), Inches(4.8),
                 font_size=Pt(14))


def build_highlight(prs, n, data):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=DARK_BG)
    add_rect(slide, 0, 0, W, Inches(0.08), fill_rgb=OX)
    footer_line(slide)
    slide_number(slide, n, TOTAL, color=RGBColor(0x88, 0x77, 0x66))

    add_text(slide, data["title"],
             Inches(0.5), Inches(0.18), Inches(11), Inches(0.7),
             font_size=Pt(34), bold=True, color=OX)
    add_text(slide, data["tagline"],
             Inches(0.5), Inches(0.88), Inches(12), Inches(0.48),
             font_size=Pt(15), color=RGBColor(0xCC, 0xBB, 0xAA))

    row_h = Inches(0.9)
    for i, (label, desc) in enumerate(data["points"]):
        y = Inches(1.55) + i * (row_h + Inches(0.12))
        add_rect(slide, Inches(0.4), y, Inches(2.6), row_h, fill_rgb=OX)
        add_text(slide, label,
                 Inches(0.5), y + Inches(0.18), Inches(2.4), row_h - Inches(0.2),
                 font_size=Pt(15), bold=True, color=WARM_WHITE)
        add_rect(slide, Inches(3.1), y, Inches(9.8), row_h,
                 fill_rgb=RGBColor(0x2A, 0x24, 0x1A))
        add_text(slide, desc,
                 Inches(3.25), y + Inches(0.2), Inches(9.5), row_h - Inches(0.2),
                 font_size=Pt(14), color=RGBColor(0xE0, 0xD4, 0xC4))


def build_feature_grid(prs, n, _data):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=WARM_WHITE)
    accent_bar(slide)
    footer_line(slide)
    slide_number(slide, n, TOTAL)

    add_rect(slide, 0, Inches(0.08), W, Inches(0.9), fill_rgb=LIGHT_SAND)
    add_text(slide, "Feature Highlights",
             Inches(0.45), Inches(0.14), Inches(11), Inches(0.6),
             font_size=Pt(30), bold=True, color=OX_DARK)

    features = [
        ("Scan",       "Point at any folder — results in seconds, no config required"),
        ("Compare",    "Side-by-side diff of any two past scans with per-file deltas"),
        ("Trend",      "SLOC growth charts over time for watched directories"),
        ("Test Metrics","Test file detection, test/code ratio, and LCOV coverage attach"),
        ("CI Gates",   "--fail-below, --fail-on-budget, --fail-above-baseline for pipelines"),
        ("Export",     "HTML · PDF · CSV · Excel · JSON — all from a single command"),
    ]

    cols = 3
    cw = Inches(3.9)
    ch = Inches(1.9)
    pad_x = Inches(0.38)
    pad_y = Inches(1.2)

    for idx, (title, desc) in enumerate(features):
        col = idx % cols
        row = idx // cols
        x = pad_x + col * (cw + Inches(0.35))
        y = pad_y + row * (ch + Inches(0.25))
        add_rect(slide, x, y, cw, ch, fill_rgb=RGBColor(0xFF, 0xFF, 0xFF),
                 line_rgb=RGBColor(0xD4, 0xC4, 0xB0), line_width=Pt(1.2))
        add_rect(slide, x, y, cw, Inches(0.36), fill_rgb=OX)
        add_text(slide, title, x + Inches(0.15), y + Inches(0.05),
                 cw - Inches(0.2), Inches(0.28),
                 font_size=Pt(13), bold=True, color=WARM_WHITE)
        add_text(slide, desc, x + Inches(0.18), y + Inches(0.45),
                 cw - Inches(0.28), ch - Inches(0.55),
                 font_size=Pt(12), color=CHARCOAL, word_wrap=True)


def build_who_benefits(prs, n):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=WARM_WHITE)
    accent_bar(slide)
    footer_line(slide)
    slide_number(slide, n, TOTAL)

    add_rect(slide, 0, Inches(0.08), W, Inches(0.9), fill_rgb=LIGHT_SAND)
    add_text(slide, "Who Benefits?",
             Inches(0.45), Inches(0.14), Inches(11), Inches(0.6),
             font_size=Pt(30), bold=True, color=OX_DARK)

    roles = [
        ("Engineering\nManagers",  OX,           ["Track scope creep sprint-over-sprint", "Inform staffing and estimation", "Provide metrics for exec reporting"]),
        ("Software\nEngineers",    OX_DARK,       ["Validate refactors reduced complexity", "Monitor test coverage growth", "Enforce code-size budgets in CI"]),
        ("Security &\nCompliance", RGBColor(0x1A, 0x4A, 0x7A), ["Auditable metrics for compliance reviews", "Detect large unexplained code growth", "Produce signed JSON artifacts for records"]),
        ("DevOps /\nPlatform",     GREEN,         ["CI gates prevent runaway growth", "Webhook-driven automated scans", "Docker-deployable for shared team use"]),
    ]

    rw = Inches(2.9)
    rh = Inches(4.6)
    start_x = Inches(0.35)
    top_y = Inches(1.15)

    for i, (role, color, points) in enumerate(roles):
        x = start_x + i * (rw + Inches(0.26))
        add_rect(slide, x, top_y, rw, Inches(0.7), fill_rgb=color)
        add_text(slide, role, x + Inches(0.12), top_y + Inches(0.06),
                 rw - Inches(0.15), Inches(0.6),
                 font_size=Pt(13), bold=True, color=WARM_WHITE)
        bullet_block(slide, points,
                     x + Inches(0.12), top_y + Inches(0.8),
                     rw - Inches(0.2), rh - Inches(0.9),
                     font_size=Pt(12), color=CHARCOAL)


def build_cadence(prs, n):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=WARM_WHITE)
    accent_bar(slide)
    footer_line(slide)
    slide_number(slide, n, TOTAL)

    add_rect(slide, 0, Inches(0.08), W, Inches(0.9), fill_rgb=LIGHT_SAND)
    add_text(slide, "How Often Should You Run It?",
             Inches(0.45), Inches(0.14), Inches(12), Inches(0.6),
             font_size=Pt(30), bold=True, color=OX_DARK)

    cadences = [
        ("On Every\nPull Request",     OX,      "Automated via CI. Catches scope creep early, before merge. Use --fail-above-baseline to block PRs that grow the codebase beyond agreed limits."),
        ("Every Sprint\n(Recommended)","#3A7A5A","A sprint-level snapshot gives managers a consistent view of how the codebase is evolving. Store each JSON result to power trend charts."),
        ("On Every\nRelease",          OX_DARK,  "Capture a permanent baseline at each release tag. Enables precise before/after comparisons when planning the next cycle."),
        ("On-Demand\n(Ad-hoc)",        RGBColor(0x1A, 0x4A, 0x7A),"Useful for audits, due-diligence reviews, vendor handoffs, or whenever a stakeholder asks 'how big is the codebase?'"),
    ]

    cw = Inches(2.9)
    ch = Inches(4.35)
    start_x = Inches(0.4)
    top_y = Inches(1.15)

    for i, (label, color, desc) in enumerate(cadences):
        if isinstance(color, str):
            r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
            c = RGBColor(r, g, b)
        else:
            c = color
        x = start_x + i * (cw + Inches(0.25))
        add_rect(slide, x, top_y, cw, Inches(0.75), fill_rgb=c)
        add_text(slide, label, x + Inches(0.12), top_y + Inches(0.06),
                 cw - Inches(0.18), Inches(0.65),
                 font_size=Pt(13), bold=True, color=WARM_WHITE)
        add_rect(slide, x, top_y + Inches(0.75), cw, ch - Inches(0.75),
                 fill_rgb=RGBColor(0xFB, 0xF5, 0xEC),
                 line_rgb=RGBColor(0xD4, 0xC4, 0xB0), line_width=Pt(1))
        add_text(slide, desc, x + Inches(0.15), top_y + Inches(0.85),
                 cw - Inches(0.25), ch - Inches(1.0),
                 font_size=Pt(12), color=CHARCOAL, word_wrap=True)


def build_getting_started(prs, n):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=DARK_BG)
    add_rect(slide, 0, 0, W, Inches(0.08), fill_rgb=OX)
    footer_line(slide)
    slide_number(slide, n, TOTAL, color=RGBColor(0x88, 0x77, 0x66))

    add_text(slide, "Getting Started",
             Inches(0.5), Inches(0.2), Inches(11), Inches(0.7),
             font_size=Pt(34), bold=True, color=OX)
    add_text(slide, "From zero to your first report in under 2 minutes",
             Inches(0.5), Inches(0.88), Inches(12), Inches(0.44),
             font_size=Pt(15), color=RGBColor(0xCC, 0xBB, 0xAA))

    steps = [
        ("1",  "Clone the repo",        "git clone https://github.com/oxide-sloc/oxide-sloc"),
        ("2",  "Launch the web UI",     "bash scripts/run.sh   →  installs on first run, opens http://127.0.0.1:4317"),
        ("3",  "Point at your codebase","Enter the path to your project and click Analyze"),
        ("4",  "Review results",        "Charts, language breakdown, and per-file table appear instantly"),
    ]

    step_h = Inches(0.85)
    for i, (num, label, cmd) in enumerate(steps):
        y = Inches(1.55) + i * (step_h + Inches(0.12))
        add_rect(slide, Inches(0.4), y, Inches(0.55), step_h, fill_rgb=OX)
        add_text(slide, num, Inches(0.4), y + Inches(0.18),
                 Inches(0.55), Inches(0.5),
                 font_size=Pt(22), bold=True, color=WARM_WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(1.05), y, Inches(11.8), step_h,
                 fill_rgb=RGBColor(0x2A, 0x24, 0x1A))
        add_text(slide, label, Inches(1.2), y + Inches(0.06),
                 Inches(4.0), Inches(0.38),
                 font_size=Pt(13), bold=True, color=OX)
        add_text(slide, cmd, Inches(5.3), y + Inches(0.07),
                 Inches(7.4), Inches(0.68),
                 font_size=Pt(12), color=RGBColor(0xC8, 0xE8, 0xC0),
                 word_wrap=True)


def build_closing(prs, n):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=DARK_BG)
    # Decorative diagonal stripe
    add_rect(slide, 0, H - Inches(2.2), W, Inches(0.08), fill_rgb=OX)

    add_text(slide, "oxide-sloc",
             Inches(0.6), Inches(1.8), Inches(10), Inches(1.4),
             font_size=Pt(72), bold=True, color=OX, align=PP_ALIGN.CENTER)
    add_text(slide, "Questions?",
             Inches(0.6), Inches(3.1), W - Inches(1.2), Inches(0.8),
             font_size=Pt(34), color=WARM_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "github.com/oxide-sloc/oxide-sloc",
             Inches(0.6), H - Inches(2.0), W - Inches(1.2), Inches(0.45),
             font_size=Pt(16), color=OX, align=PP_ALIGN.CENTER)
    add_text(slide, "AGPL-3.0-or-later  ·  Author: Nima Shafie  ·  nimzshafie@gmail.com",
             Inches(0.6), H - Inches(1.5), W - Inches(1.2), Inches(0.38),
             font_size=Pt(11), color=RGBColor(0x88, 0x77, 0x66), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Build all slides
# ─────────────────────────────────────────────────────────────────────────────

for idx, data in enumerate(slides_data):
    t = data["type"]
    if t == "cover":
        build_cover(prs)
    elif t == "content":
        build_content(prs, idx + 1, data)
    elif t == "two_col":
        build_two_col(prs, idx + 1, data)
    elif t == "highlight":
        build_highlight(prs, idx + 1, data)
    elif t == "feature_grid":
        build_feature_grid(prs, idx + 1, data)
    elif t == "who_benefits":
        build_who_benefits(prs, idx + 1)
    elif t == "cadence":
        build_cadence(prs, idx + 1)
    elif t == "getting_started":
        build_getting_started(prs, idx + 1)
    elif t == "closing":
        build_closing(prs, idx + 1)

# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────

out_path = os.path.join(os.path.dirname(__file__), "oxide-sloc-overview.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
